"""
Document scraper module for extracting text from various file types
"""
import os
from typing import List, Dict, Any, Optional
import io
import pytesseract
from PIL import Image
import docx2txt
import pandas as pd
from pptx import Presentation
from pdf_scraper import PDFScraper

class DocumentScraper:
    """Class for extracting text from various document types"""
    
    def __init__(self):
        """Initialize the document scraper"""
        self.pdf_scraper = PDFScraper()
        # Set pytesseract path
        pytesseract.pytesseract.tesseract_cmd = r'C:\\Users\\DavidGidony\\AppData\\Local\\Programs\\Tesseract-OCR\\tesseract.exe'
    
    def extract_text(self, file_path: str) -> str:
        """
        Extract text from a file based on its extension
        
        Args:
            file_path: Path to the file
            
        Returns:
            Extracted text from the file
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_ext == '.pdf':
                return self.pdf_scraper.extract_text_from_pdf(file_path)
            elif file_ext == '.txt':
                return self._extract_text_from_txt(file_path)
            elif file_ext in ['.docx', '.doc']:
                return self._extract_text_from_word(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                return self._extract_text_from_excel(file_path)
            elif file_ext in ['.pptx', '.ppt']:
                return self._extract_text_from_powerpoint(file_path)
            elif file_ext in ['.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff']:
                return self._extract_text_from_image(file_path)
            else:
                return f"Unsupported file type: {file_ext}"
        except Exception as e:
            raise Exception(f"Error extracting text from {file_path}: {str(e)}")
    
    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        """
        Extract metadata from a file based on its extension
        
        Args:
            file_path: Path to the file
            
        Returns:
            Dictionary containing metadata
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_ext = os.path.splitext(file_path)[1].lower()
        file_name = os.path.basename(file_path)
        file_size = os.path.getsize(file_path)
        
        metadata = {
            "file_name": file_name,
            "file_path": file_path,
            "file_type": file_ext,
            "file_size": file_size,
        }
        
        try:
            if file_ext == '.pdf':
                pdf_metadata = self.pdf_scraper.extract_metadata(file_path)
                metadata.update(pdf_metadata)
        except Exception:
            pass
        
        return metadata
    
    def _extract_text_from_txt(self, file_path: str) -> str:
        """Extract text from a TXT file"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            return file.read()
    
    def _extract_text_from_word(self, file_path: str) -> str:
        """Extract text from a Word document"""
        return docx2txt.process(file_path)
    
    def _extract_text_from_excel(self, file_path: str) -> str:
        """Extract text from an Excel file"""
        df = pd.read_excel(file_path, sheet_name=None)
        text = ""
        for sheet_name, sheet_df in df.items():
            text += f"Sheet: {sheet_name}\n"
            text += sheet_df.to_string(index=False) + "\n\n"
        return text
    
    def _extract_text_from_powerpoint(self, file_path: str) -> str:
        """Extract text from a PowerPoint presentation"""
        prs = Presentation(file_path)
        text = ""
        for i, slide in enumerate(prs.slides):
            text += f"Slide {i+1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"
        return text
    
    def _extract_text_from_image(self, file_path: str) -> str:
        """Extract text from an image using OCR"""
        img = Image.open(file_path)
        return pytesseract.image_to_string(img)
    
    def chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """
        Split text into chunks with overlap
        
        Args:
            text: Text to split
            chunk_size: Size of each chunk
            overlap: Overlap between chunks
            
        Returns:
            List of text chunks
        """
        return self.pdf_scraper.chunk_text(text, chunk_size, overlap)